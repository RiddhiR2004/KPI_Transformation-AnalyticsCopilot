from __future__ import annotations

import io
import csv
import logging
from pathlib import Path

logger = logging.getLogger("app.services.asset_parser")

def extract_text_from_asset(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()
    
    if ext == ".txt":
        return content.decode("utf-8", errors="ignore")
    
    elif ext == ".csv":
        try:
            text_lines = []
            decoded = content.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(decoded))
            for row in reader:
                if row:
                    text_lines.append(", ".join(row))
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"Error parsing CSV file {filename}: {e}")
            return f"[Error parsing CSV: {e}]"
            
    elif ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text_lines = []
            for sheet in wb.worksheets:
                text_lines.append(f"--- Spreadsheet Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_str = ", ".join([str(val) for val in row if val is not None])
                    if row_str.strip():
                        text_lines.append(row_str)
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"Error parsing XLSX file {filename}: {e}")
            return f"[Excel content parsed from {filename} with error: {e}]"

    elif ext == ".xls":
        return f"[Legacy Excel spreadsheet: {filename}. Upgrade to .xlsx for structured cell parsing.]"
            
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(io.BytesIO(content))
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error parsing DOCX file {filename}: {e}")
            return f"[Error parsing DOCX: {e}]"
            
    elif ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            logger.error(f"Error parsing PDF file {filename}: {e}")
            return f"[Error parsing PDF: {e}]"

    elif ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(content))
            text_runs = []
            for i, slide in enumerate(prs.slides):
                text_runs.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logger.error(f"Error parsing PPTX file {filename}: {e}")
            return f"[Error parsing PPTX: {e}]"

    elif ext in [".png", ".jpg", ".jpeg"]:
        return f"[Image Asset: {filename}. Multimodal/vision analysis will be used directly on this visual asset.]"
        
    else:
        raise ValueError(f"Unsupported file format: {ext}")
